import streamlit as st
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
from langchain_text_splitters import RecursiveCharacterTextSplitter
import io

# -------------------------
# Page Config
# -------------------------
st.set_page_config(page_title="Meeting Intelligence RAG", layout="wide")
st.title("Meeting Intelligence RAG System")
st.write("Upload a meeting transcript and get AI-powered summaries, decisions, action items, and Q&A.")

# -------------------------
# File Parsers
# -------------------------

def extract_text_from_txt(file):
    return file.read().decode("utf-8")

def extract_text_from_pdf(file):
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(file.read()))
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text.strip()
    except Exception as e:
        st.error(f"Error reading PDF: {e}")
        return ""

def extract_text_from_pptx(file):
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file.read()))
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text.strip() + "\n"
        return text.strip()
    except Exception as e:
        st.error(f"Error reading PPTX: {e}")
        return ""

# -------------------------
# Load Models (cached)
# -------------------------

@st.cache_resource
def load_embedding_model():
    return SentenceTransformer("all-MiniLM-L6-v2")

@st.cache_resource
def load_llm():
    model_name = "google/flan-t5-large"
    tokenizer = AutoTokenizer.from_pretrained(model_name)
    model = AutoModelForSeq2SeqLM.from_pretrained(model_name)
    return tokenizer, model

# -------------------------
# LLM Helper — safely handles token limit
# -------------------------

def run_llm(tokenizer, model, prompt, max_new_tokens=150):
    inputs = tokenizer(
        prompt,
        return_tensors="pt",
        truncation=True,
        max_length=512
    )
    outputs = model.generate(
        **inputs,
        max_new_tokens=max_new_tokens,
        num_beams=4,
        early_stopping=True,
        no_repeat_ngram_size=3
    )
    return tokenizer.decode(outputs[0], skip_special_tokens=True).strip()

# -------------------------
# Smart context: get top chunks for a query
# -------------------------

def get_top_chunks(query, embed_model, index, chunks, k=5):
    q_vec = embed_model.encode([query])
    distances, indices = index.search(np.array(q_vec), k=k)
    return " ".join([chunks[i] for i in indices[0] if i < len(chunks)])

# -------------------------
# File Upload
# -------------------------

uploaded_file = st.file_uploader(
    "Upload Transcript (TXT, PDF, or PPTX)",
    type=["txt", "pdf", "pptx"]
)

if uploaded_file:
    file_type = uploaded_file.name.split(".")[-1].lower()

    with st.spinner("Reading file..."):
        if file_type == "txt":
            text = extract_text_from_txt(uploaded_file)
        elif file_type == "pdf":
            text = extract_text_from_pdf(uploaded_file)
        elif file_type == "pptx":
            text = extract_text_from_pptx(uploaded_file)
        else:
            st.error("Unsupported file type.")
            st.stop()

    if not text:
        st.error("Could not extract any text from the file. Please check the file.")
        st.stop()

    st.subheader("Transcript Preview")
    st.write(text[:1000] + ("..." if len(text) > 1000 else ""))
    st.caption(f"Total characters extracted: {len(text)}")

    # -------------------------
    # Chunking
    # -------------------------
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=400,
        chunk_overlap=80
    )
    chunks = splitter.split_text(text)

    # -------------------------
    # Load models
    # -------------------------
    with st.spinner("Loading AI models (first load takes ~1 min)..."):
        embed_model = load_embedding_model()
        tokenizer, llm_model = load_llm()

    # -------------------------
    # Build FAISS index
    # -------------------------
    embeddings = embed_model.encode(chunks)
    dimension = embeddings.shape[1]
    index = faiss.IndexFlatL2(dimension)
    index.add(np.array(embeddings))

    # -------------------------
    # Meeting Summary
    # — use top chunks relevant to "summary overview"
    # -------------------------
    st.subheader("Meeting Summary")
    with st.spinner("Generating summary..."):
        summary_context = get_top_chunks(
            "main topic discussion agenda overview", embed_model, index, chunks, k=6
        )
        summary_prompt = (
            "Read the meeting notes below and write a clear, concise summary "
            "of what the meeting was about in 3-4 sentences.\n\n"
            f"Meeting notes:\n{summary_context}\n\n"
            "Summary:"
        )
        summary = run_llm(tokenizer, llm_model, summary_prompt, max_new_tokens=150)
    st.write(summary)

    # -------------------------
    # Key Decisions
    # -------------------------
    st.subheader("Key Decisions")
    with st.spinner("Extracting decisions..."):
        decision_context = get_top_chunks(
            "decision agreed decided approved confirmed concluded", embed_model, index, chunks, k=5
        )
        decision_prompt = (
            "From the meeting notes below, list the key decisions that were made. "
            "Each decision should be on a new line starting with a dash.\n\n"
            f"Meeting notes:\n{decision_context}\n\n"
            "Key decisions made:\n-"
        )
        decisions = run_llm(tokenizer, llm_model, decision_prompt, max_new_tokens=150)
    st.write("- " + decisions if not decisions.startswith("-") else decisions)

    # -------------------------
    # Action Items
    # -------------------------
    st.subheader("Action Items")
    with st.spinner("Extracting action items..."):
        action_context = get_top_chunks(
            "task assigned will do responsible person action item follow up", embed_model, index, chunks, k=5
        )
        action_prompt = (
            "From the meeting notes below, extract specific tasks assigned to people. "
            "Format each task as: Person Name - Task description\n\n"
            f"Meeting notes:\n{action_context}\n\n"
            "Action items:\n"
        )
        actions = run_llm(tokenizer, llm_model, action_prompt, max_new_tokens=150)
    st.write(actions)

    # -------------------------
    # Question Answering
    # -------------------------
    st.subheader("Ask a Question About the Meeting")
    question = st.text_input("Type your question here and press Enter")

    if question and question.strip():
        with st.spinner("Searching transcript and generating answer..."):
            context = get_top_chunks(question, embed_model, index, chunks, k=5)
            qa_prompt = (
                "You are a helpful assistant. Answer the question below using "
                "only the provided meeting context. Be specific and factual.\n\n"
                f"Context:\n{context}\n\n"
                f"Question: {question}\n\n"
                "Answer:"
            )
            answer = run_llm(tokenizer, llm_model, qa_prompt, max_new_tokens=150)

        st.subheader("Answer")
        st.success(answer)

        if not answer or len(answer) < 5:
            st.warning(
                "The model could not find a confident answer. "
                "Try rephrasing your question or check if the information is in the transcript."
            )
